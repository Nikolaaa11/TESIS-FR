"""
exportar_ppt.py
---------------
Genera docs/Tesis_USS_Defensa.pptx: presentación de defensa de tesis (16:9) con
marca Universidad San Sebastián, hallazgos, figuras y diapositiva de objeciones.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.append(os.path.dirname(__file__))
import config as C

NAVY = RGBColor(0x16, 0x24, 0x3A)
COPPER = RGBColor(0xC2, 0x70, 0x3D)
GRAY = RGBColor(0x5A, 0x5A, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
INK = RGBColor(0x1D, 0x1D, 0x1F)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FIG = C.FIG


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    r.font.italic = italic
    return tb


def _bar(slide, color=COPPER, y=Inches(1.55), h=Pt(3)):
    sh = slide.shapes.add_shape(1, Inches(0.7), y, Inches(2.2), h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()


def add_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, NAVY)
    _box(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.6),
         "UNIVERSIDAD SAN SEBASTIÁN · Magíster en Data Science", 18, COPPER, True)
    _box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.4),
         "El cobre, la bolsa y la liquidez", 54, WHITE, True, font="Georgia")
    _box(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(1.4),
         "Impacto de las variables macroeconómicas globales y financieras en la "
         "valoración bursátil del sector de minería de cobre en Chile (2004–2026)",
         22, RGBColor(0xCF, 0xD6, 0xE2))
    _box(s, Inches(0.9), Inches(6.3), Inches(11), Inches(0.8),
         "Defensa de tesis · Autor: ____________  ·  Profesor guía: ____________  ·  2026",
         14, RGBColor(0xAE, 0xB4, 0xC2))


def add_section(prs, kicker, title, bullets, foot=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _box(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.4), kicker, 14, COPPER, True)
    _box(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.9), title, 32, NAVY, True, font="Georgia")
    _bar(s)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl0 = not b.startswith("  ")
        r = p.add_run(); r.text = ("•  " if lvl0 else "–  ") + b.strip()
        r.font.size = Pt(19 if lvl0 else 16); r.font.color.rgb = INK if lvl0 else GRAY
        r.font.name = "Calibri"; p.space_after = Pt(9); p.level = 0 if lvl0 else 1
    if foot:
        _box(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4), foot, 12, GRAY, italic=True)


def add_image(prs, kicker, title, img, caption):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _box(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.4), kicker, 14, COPPER, True)
    _box(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.9), title, 30, NAVY, True, font="Georgia")
    _bar(s)
    p = FIG / img
    if p.exists():
        try:
            from PIL import Image as PImg
            iw, ih = PImg.open(str(p)).size
            maxw, maxh = Inches(8.6), Inches(4.8)
            w = maxw; h = Emu(int(w * ih / iw))
            if h > maxh:
                h = maxh; w = Emu(int(h * iw / ih))
            s.shapes.add_picture(str(p), Inches(0.7), Inches(2.0), width=w, height=h)
        except Exception:
            pass
    _box(s, Inches(9.5), Inches(2.2), Inches(3.3), Inches(4.5), caption, 16, GRAY)


def add_table(prs, kicker, title, headers, rows, caption=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _box(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.4), kicker, 14, COPPER, True)
    _box(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.9), title, 30, NAVY, True, font="Georgia")
    _bar(s)
    nr, nc = len(rows) + 1, len(headers)
    tbl = s.shapes.add_table(nr, nc, Inches(0.7), Inches(2.0), Inches(12), Inches(0.4 * nr)).table
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j); c.text = htxt
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(14); pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            rp = c.text_frame.paragraphs[0]; rp.runs[0].font.size = Pt(13); rp.runs[0].font.color.rgb = INK
    if caption:
        _box(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4), caption, 12, GRAY, italic=True)


def add_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, NAVY)
    _box(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6),
         "La iliquidez retrasa, pero no elimina,\nla transmisión del cobre", 40, WHITE, True, font="Georgia")
    _box(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.6),
         "Gracias.  ·  github.com/Nikolaaa11/TESIS-FR  ·  web-pi-pied-45.vercel.app",
         16, RGBColor(0xCF, 0xD6, 0xE2))


def build():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    add_title(prs)
    add_section(prs, "01 · El problema", "Una pregunta poco estudiada", [
        "Chile: ~¼ de la producción mundial de cobre; el metal es ~½ de las exportaciones.",
        "Brecha: la literatura cubre cobre→tipo de cambio→macro, pero no cobre→valoración minera.",
        "Restricción: universo de mineras de cobre cotizadas mínimo (Codelco no cotiza).",
        "Pregunta: magnitud, dinámica y rol de la liquidez en la transmisión cobre→acción.",
    ])
    add_section(prs, "02 · Universo", "Resuelto con evidencia empírica", [
        "Antofagasta plc (LSE, GBp): pure-play líquido, grupo Luksic, ~654 kt de cobre (2025).",
        "Pucobre (Santiago, CLP): único pure-play de cobre listado en Chile; ilíquido.",
        "  62% de días sin transacción — el contraste clave de la tesis.",
        "Panel materiales (CAP, SQM) y referencias internacionales (SCCO, FCX, BHP, GLEN).",
    ])
    add_section(prs, "03 · Método", "El modelo nace de los datos", [
        "Log-precios I(1), retornos I(0) (ADF/PP/KPSS + Zivot-Andrews).",
        "Corto plazo: regresión HAC, VAR/IRF/FEVD, causalidad de Toda-Yamamoto.",
        "Largo plazo: cointegración de Johansen, VECM, ARDL/NARDL bounds.",
        "Riesgo y micro: GARCH/EGARCH/GJR, panel FE/RE, event study, iliquidez de Amihud.",
        "Identificación: cobre como shock exógeno + cuasi-experimento líquido vs ilíquido.",
    ])
    add_section(prs, "04 · Hallazgo central", "La transmisión crece con el horizonte", [
        "Antofagasta (líquido): elasticidad-cobre 0.70; el cobre explica ~28% de su varianza.",
        "Pucobre (ilíquido): contemporánea ≈ 0.09 (R²=0.04), pero crece con el horizonte:",
        "  0.09 diaria  →  0.42 acumulada 5d  →  0.60 mensual  →  0.75 largo plazo.",
        "El 95% del impacto llega el día 0 en ANTO; sólo el 29% en Pucobre.",
        "La iliquidez retrasa, pero no elimina, el vínculo fundamental cobre→valoración.",
    ], foot="Triangulado por 6 métodos independientes.")
    add_image(prs, "04 · Hallazgo central", "Curva de transmisión y ciclo del cobre",
              "precios_normalizados.png",
              "Los pure-plays amplifican el ciclo del cobre (apalancamiento operativo). "
              "Pucobre lo sigue con rezago por su iliquidez.")
    add_table(prs, "05 · Resultados", "Elasticidad-cobre contemporánea (HAC)",
              ["Activo", "β-cobre", "t", "R²"],
              [["Antofagasta", "0.70", "15.4", "0.42"], ["FCX", "0.63", "10.3", "0.53"],
               ["Glencore", "0.58", "9.8", "0.29"], ["Southern", "0.49", "9.6", "0.58"],
               ["BHP", "0.32", "10.7", "0.60"], ["Pucobre", "0.09", "4.4", "0.04"]],
              caption="Pucobre es el atípico de baja transmisión contemporánea pese a ser pure-play.")
    add_image(prs, "05 · Resultados", "Impulso-respuesta: respuesta diferida de Pucobre",
              "irf_PUCOBRE_SN.png",
              "La respuesta de Pucobre al shock de cobre se ACUMULA en los días siguientes "
              "(IRF 5d ≈ 4× el día 1), no contemporáneamente.")
    add_table(prs, "06 · Largo plazo", "Cointegración (Johansen r=1) y VECM",
              ["Activo", "Elast. cobre LP", "Ajuste α"],
              [["Antofagasta", "0.86", "−0.0008"], ["Pucobre", "0.75", "−0.0004"]],
              caption="En el largo plazo Pucobre (0.75) es comparable a Antofagasta (0.86).")
    add_image(prs, "07 · Iliquidez", "Iliquidez vs transmisión contemporánea",
              "iliquidez_vs_beta.png",
              "Relación negativa robusta a 4 proxies (Amihud, % ceros, Roll, volumen). "
              "Pucobre: 62% de días sin transar.")
    add_section(prs, "08 · Robustez", "Pruebas adicionales de rigor", [
        "Quiebre estructural (Quandt-Andrews): ANTO duplica su β en 2007 (0.34→0.78);",
        "  Pucobre sin quiebre — su débil transmisión es estructuralmente estable.",
        "Iliquidez multi-proxy: 4 medidas distintas, mismo signo (robusto).",
        "Fuera de muestra (Clark-West): el cobre rezagado predice Pucobre (p=0.004),",
        "  más que ANTO; placebo: SQM (litio) no responde (p=0.68).",
        "Volatilidad: persistencia ≈0.99 y efecto apalancamiento (GJR/EGARCH).",
        "NARDL: asimetría de largo plazo significativa en Antofagasta (p=0.004).",
    ])
    add_table(prs, "09 · Síntesis", "Verificación de hipótesis",
              ["H", "Enunciado", "Veredicto"],
              [["H1", "Cobre positivo y significativo", "Se sostiene"],
               ["H2", "Sensibilidad según liquidez", "Se sostiene"],
               ["H3", "Cointegración de largo plazo", "Se sostiene"],
               ["H4", "Volatilidad persistente/asimétrica", "Se sostiene"],
               ["H5", "Transmisión diferida (iliquidez)", "Se sostiene con fuerza"]])
    add_section(prs, "10 · Conclusiones", "Aportes e implicancias", [
        "Contribución: curva de transmisión cobre→valoración por horizonte, atribuida a microestructura.",
        "Para inversionistas: una beta diaria subestima la exposición de un ilíquido ~8×.",
        "Para el mercado: cuantifica una ineficiencia de corto plazo por liquidez, no por fundamento.",
        "Líneas futuras: EMBI/BCCh, NARDL dinámico, DCC-GARCH, sorpresa de TPM.",
    ])
    add_section(prs, "11 · Defensa", "Cinco objeciones y respuestas", [
        "Endogeneidad: Toda-Yamamoto unidireccional cobre→acción; el cobre es global y exógeno.",
        "Quiebres en 22 años: testeados (Quandt-Andrews); el hallazgo es estable.",
        "N pequeño del panel: se reconoce; la inferencia primaria es por serie de tiempo, no panel.",
        "¿Depende de Amihud?: no; robusto a 4 proxies de iliquidez.",
        "Validez externa: confirmado por benchmarks internacionales y placebo (SQM).",
    ])
    add_closing(prs)

    out = C.ROOT / "docs" / "Tesis_USS_Defensa.pptx"
    try:
        prs.save(out); print(f"OK {out} ({len(prs.slides._sldIdLst)} slides)")
    except PermissionError:
        out = C.ROOT / "docs" / "Tesis_USS_Defensa_v2.pptx"; prs.save(out)
        print(f"AVISO bloqueado; guardado como {out}")
    # copia a web
    web = C.ROOT / "web" / "assets" / "docs" / "Tesis_USS_Defensa.pptx"
    web.parent.mkdir(parents=True, exist_ok=True)
    import shutil
    try:
        shutil.copy(out, web); print(f"OK {web}")
    except Exception as e:
        print("copia web fallo:", e)


if __name__ == "__main__":
    build()
